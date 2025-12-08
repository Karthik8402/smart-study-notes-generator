"""File Extraction Module"""

import os
import re
from typing import Tuple
import asyncio


async def extract_pdf_text(file_path: str) -> str:
    """Extract text from PDF file."""
    from pypdf import PdfReader
    
    def _read_pdf():
        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n\n".join(text_parts)
    
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _read_pdf)


async def extract_ppt_text(file_path: str) -> str:
    """Extract text from PowerPoint file."""
    from pptx import Presentation
    
    def _read_ppt():
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _read_ppt)


async def extract_image_text(file_path: str) -> str:
    """Extract text from image using OCR."""
    try:
        import pytesseract
        from PIL import Image
        
        def _ocr():
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, _ocr)
    except Exception as e:
        return f"OCR Error: {str(e)}"


async def extract_plain_text(file_path: str) -> str:
    """Extract text from plain text file."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


async def extract_youtube_transcript(url: str) -> Tuple[str, str]:
    """Extract transcript from YouTube video."""
    from youtube_transcript_api import YouTubeTranscriptApi
    
    # Extract video ID from various URL formats
    video_id = None
    
    # Pattern for youtu.be/VIDEO_ID (with optional query params)
    if 'youtu.be/' in url:
        match = re.search(r'youtu\.be\/([a-zA-Z0-9_-]{11})', url)
        if match:
            video_id = match.group(1)
    
    # Pattern for youtube.com/watch?v=VIDEO_ID
    if not video_id and 'v=' in url:
        match = re.search(r'v=([a-zA-Z0-9_-]{11})', url)
        if match:
            video_id = match.group(1)
    
    # Pattern for youtube.com/embed/VIDEO_ID or youtube.com/shorts/VIDEO_ID
    if not video_id:
        match = re.search(r'(?:embed|shorts)\/([a-zA-Z0-9_-]{11})', url)
        if match:
            video_id = match.group(1)
    
    # Direct video ID (11 characters)
    if not video_id and len(url.strip()) == 11:
        video_id = url.strip()
    
    # Error if no video ID found
    if not video_id:
        raise ValueError("Could not extract YouTube video ID from URL")
    
    def _get_transcript():
        try:
            # Use the correct API - it's a class method
            ytt_api = YouTubeTranscriptApi()
            transcript_list = ytt_api.fetch(video_id)
            text = " ".join([entry.text for entry in transcript_list])
            return text
        except Exception as e:
            # Fallback: try older API format
            try:
                transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
                text = " ".join([entry['text'] for entry in transcript_list])
                return text
            except:
                raise ValueError(f"Could not fetch transcript: {str(e)}")
    
    loop = asyncio.get_event_loop()
    text = await loop.run_in_executor(None, _get_transcript)
    
    if not text or len(text) < 50:
        raise ValueError("Transcript too short or empty")
    
    return text, f"YouTube Video {video_id}"


def clean_extracted_text(text: str) -> str:
    """Clean and normalize extracted text."""
    if not text:
        return ""
    
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Remove special characters but keep basic punctuation
    text = re.sub(r'[^\w\s.,!?;:\-\'"()\[\]{}@#$%&*+=/]', '', text)
    
    # Normalize line breaks
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    return text.strip()
